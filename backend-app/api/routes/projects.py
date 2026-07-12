from fastapi import APIRouter, HTTPException, Depends, UploadFile, File
from sqlalchemy.orm import Session
from loguru import logger
from langchain_core.messages import HumanMessage
import uuid
import datetime

from models.schemas import (
    PromptRequest, 
    InterpretationResponse, 
    ProjectInterpretationDTO, 
    ExecuteRequest,
    ProjectPlan,
    SavePlanRequest,
    GanttTaskUpdate,
    GanttLinkCreate,
    GanttLinkUpdate
)
from core.database import get_db
from models.domain import Project, Task, Subtask, TaskLink
from services.ai.agents import interpreter_agent, planner_agent
from core.date_helper import WeekendHelper

router = APIRouter()

@router.post("/extract-file")
async def extract_file_content(file: UploadFile = File(...)):
    """
    Extracts text from the uploaded file (PDF, DOCX, PPTX, XLSX, TXT)
    """
    content = await file.read()
    filename = file.filename.lower()
    text = ""
    
    try:
        if filename.endswith('.txt') or filename.endswith('.md'):
            text = content.decode('utf-8')
        elif filename.endswith('.pdf'):
            import PyPDF2
            import io
            pdf_reader = PyPDF2.PdfReader(io.BytesIO(content))
            for page in pdf_reader.pages:
                text += page.extract_text() + "\n"
        elif filename.endswith('.docx'):
            import docx
            import io
            doc = docx.Document(io.BytesIO(content))
            text = "\n".join([para.text for para in doc.paragraphs])
        elif filename.endswith('.pptx'):
            from pptx import Presentation
            import io
            prs = Presentation(io.BytesIO(content))
            for slide in prs.slides:
                for shape in slide.shapes:
                    if hasattr(shape, "text"):
                        text += shape.text + "\n"
        elif filename.endswith('.xlsx'):
            import openpyxl
            import io
            wb = openpyxl.load_workbook(io.BytesIO(content), data_only=True)
            for sheet in wb.worksheets:
                for row in sheet.iter_rows(values_only=True):
                    row_text = [str(cell) for cell in row if cell is not None]
                    if row_text:
                        text += " | ".join(row_text) + "\n"
        else:
            raise HTTPException(status_code=400, detail="Unsupported file format")
            
        return {"extractedText": text.strip()}
    except Exception as e:
        logger.error(f"Error extracting text from {filename}: {str(e)}")
        raise HTTPException(status_code=500, detail="Error extracting file content")


@router.post("/interpret", response_model=InterpretationResponse)
async def interpret_prompt(request: PromptRequest):
    """
    Interprets a natural language prompt and returns a structured DTO (Pydantic Model)
    representing the intended action. Does not modify the database.
    """
    logger.info(f"Received prompt for interpretation: {request.prompt}")
    
    try:
        # Pass prompt to the Interpreter Agent
        state = {
            "messages": [HumanMessage(content=request.prompt)],
            "context": "",
            "parsed_output": {}
        }
        
        result = interpreter_agent.invoke(state)
        parsed_data = result.get("parsed_output", {})
        
        if not parsed_data:
            logger.warning("Agent could not parse a valid structured output.")
            raise HTTPException(status_code=400, detail="Could not interpret the prompt.")
        
        dto = ProjectInterpretationDTO(**parsed_data)
        
        # Determine action based on clarification agent's missing fields
        missing = parsed_data.get("missingFields", [])
        action = "NEEDS_CLARIFICATION" if missing else "READY_TO_PLAN"
        
        # Parse confidence
        conf_str = parsed_data.get("planningConfidence", "0%")
        conf_val = float(conf_str.strip('%')) / 100.0 if conf_str.strip('%').isdigit() else 0.0
        
        warnings = []
        if missing:
            warnings.append(f"Missing required fields: {', '.join(missing)}")
        
        response = InterpretationResponse(
            action=action,
            confidence=conf_val,
            dto=dto,
            warnings=warnings,
            requiresConfirmation=True
        )
        
        return response
        
    except Exception as e:
        logger.error(f"Error interpreting prompt: {str(e)}")
        raise HTTPException(status_code=500, detail="Internal Server Error during interpretation.")

@router.post("/execute", response_model=ProjectPlan)
async def execute_action(payload: ExecuteRequest):
    """
    Executes a confirmed action. 
    Takes the confirmed DTO and passes it to the Planner Agent to generate a full
    detailed Jira-style project breakdown.
    """
    logger.info(f"Executing confirmed action and generating plan for: {payload.confirmedDto.projectTitle}")
    
    try:
        # Pass confirmed DTO to the Planner Agent
        state = {
            "project_dto": payload.confirmedDto.model_dump(),
            "project_plan": {}
        }
        
        result = planner_agent.invoke(state)
        plan_data = result.get("project_plan", {})
        
        if not plan_data:
            logger.warning("Planner Agent failed to generate a project plan.")
            raise HTTPException(status_code=500, detail="Failed to generate project plan.")
            
        # Return the generated Project Plan
        plan = ProjectPlan(**plan_data)
        return plan
        
    except Exception as e:
        logger.error(f"Error generating project plan: {str(e)}")
        raise HTTPException(status_code=500, detail="Internal Server Error during execution.")


@router.post("/save")
async def save_project_plan(payload: SavePlanRequest, db: Session = Depends(get_db)):
    """
    Saves the finalized project plan into the PostgreSQL database.
    """
    logger.info(f"Saving project plan: {payload.dto.projectTitle}")
    try:
        dto = payload.dto
        plan = payload.plan
        
        # 1. Create Project
        project_code = f"PRJ-{str(uuid.uuid4())[:8].upper()}"
        
        planned_start = None
        if dto.plannedStartDate:
            try:
                # Try to parse standard ISO format
                from datetime import datetime
                planned_start = datetime.strptime(dto.plannedStartDate, "%Y-%m-%d").date()
            except Exception:
                pass # skip if LLM gave conversational date
                
        db_project = Project(
            project_code=project_code,
            project_title=dto.projectTitle or "Untitled Project",
            project_type=dto.projectType or "General",
            vessel_name=dto.vesselName or "Unknown",
            vessel_type=dto.vesselType or "Unknown",
            scope_summary=dto.scopeSummary,
            priority_level=dto.priorityLevel,
            planned_start_date=planned_start,
            duration_weeks=dto.durationWeeks,
            budget_at_completion=dto.budgetAtCompletion,
            currency=dto.currency,
            crew_size=dto.crewSize,
            known_risks=", ".join(dto.knownRisks) if dto.knownRisks else None,
            weather_constraints=dto.weatherConstraints,
            status="To Do"
        )
        db.add(db_project)
        db.flush() # to get db_project.project_id
        
        # 2. Create Tasks & Subtasks
        for t_idx, task_dto in enumerate(plan.tasks):
            task_code = f"{project_code}-T{t_idx+1}"
            
            t_planned_start = None
            t_planned_end = None
            if getattr(task_dto, "planned_start_date", None):
                try:
                    t_planned_start = datetime.strptime(task_dto.planned_start_date, "%Y-%m-%d").date()
                except ValueError:
                    pass
            if getattr(task_dto, "planned_end_date", None):
                try:
                    t_planned_end = datetime.strptime(task_dto.planned_end_date, "%Y-%m-%d").date()
                except ValueError:
                    pass
                    
            db_task = Task(
                task_code=task_code,
                project_id=db_project.project_id,
                task_name=task_dto.summary,
                task_description=task_dto.description,
                priority=task_dto.priority,
                status="To Do",
                planned_start_date=t_planned_start,
                planned_end_date=t_planned_end,
                actual_start_date=t_planned_start,
                actual_end_date=t_planned_end,
                planned_cost=task_dto.planned_cost or 0.0,
                actual_cost=task_dto.actual_cost or task_dto.planned_cost or 0.0,
                planned_crew=getattr(task_dto, "planned_crew", None) or getattr(task_dto, "crew_needed", None) or 1,
                actual_crew=getattr(task_dto, "actual_crew", None) or getattr(task_dto, "planned_crew", None) or getattr(task_dto, "crew_needed", None) or 1,
                order_index=task_dto.order_index or (t_idx + 1)
            )
            db.add(db_task)
            db.flush()
            
            for s_idx, sub_dto in enumerate(task_dto.subtasks):
                sub_code = f"{task_code}-S{s_idx+1}"
                
                s_planned_start = None
                s_planned_end = None
                if getattr(sub_dto, "planned_start_date", None):
                    try:
                        s_planned_start = datetime.strptime(sub_dto.planned_start_date, "%Y-%m-%d").date()
                    except ValueError:
                        pass
                if getattr(sub_dto, "planned_end_date", None):
                    try:
                        s_planned_end = datetime.strptime(sub_dto.planned_end_date, "%Y-%m-%d").date()
                    except ValueError:
                        pass
                        
                db_subtask = Subtask(
                    subtask_code=sub_code,
                    task_id=db_task.task_id,
                    subtask_name=sub_dto.summary,
                    subtask_description=sub_dto.description,
                    planned_start_date=s_planned_start,
                    planned_end_date=s_planned_end,
                    actual_start_date=s_planned_start,
                    actual_end_date=s_planned_end,
                    status="To Do",
                    planned_cost=sub_dto.planned_cost or 0.0,
                    actual_cost=sub_dto.actual_cost or sub_dto.planned_cost or 0.0,
                    planned_crew=getattr(sub_dto, "planned_crew", None) or getattr(sub_dto, "crew_needed", None) or 1,
                    actual_crew=getattr(sub_dto, "actual_crew", None) or getattr(sub_dto, "planned_crew", None) or getattr(sub_dto, "crew_needed", None) or 1,
                    order_index=sub_dto.order_index or (s_idx + 1)
                )
                db.add(db_subtask)
                
        db.commit()
        return {"status": "success", "message": "Project plan saved successfully", "project_id": str(db_project.project_id)}
    except Exception as e:
        db.rollback()
        logger.error(f"Error saving project plan to DB: {str(e)}")
        raise HTTPException(status_code=500, detail=str(e))

@router.get("/")
def get_all_projects(db: Session = Depends(get_db)):
    """
    Returns a list of all projects.
    """
    try:
        projects = db.query(Project).all()
        projects_list = [{
            "project_id": str(p.project_id),
            "project_code": p.project_code,
            "project_title": p.project_title,
            "project_type": p.project_type,
            "vessel_name": p.vessel_name,
            "status": p.status,
            "planned_start_date": p.planned_start_date.isoformat() if p.planned_start_date else None,
            "duration_weeks": p.duration_weeks,
            "currency": p.currency or "USD"
        } for p in projects]
        
        # Add demo project
        projects_list.insert(0, {
            "project_id": "demo-maritime-project-001",
            "project_code": "DEMO-001",
            "project_title": "Vessel MV-Alpha Retrofit (DEMO)",
            "project_type": "Retrofit",
            "vessel_name": "MV-Alpha",
            "status": "In Progress",
            "planned_start_date": "2019-05-31",
            "duration_weeks": 208, # ~4 years
            "currency": "EUR"
        })
        
        return projects_list
    except Exception as e:
        logger.error(f"Error fetching projects: {str(e)}")
        raise HTTPException(status_code=500, detail="Internal Server Error")

@router.get("/{project_id}/gantt")
def get_project_gantt_data(project_id: str, db: Session = Depends(get_db)):
    """
    Returns data formatted for DHTMLX Gantt for a specific project.
    """
    try:
        if project_id == "demo-maritime-project-001":
            demo_data = {
                "data": [
                    { "id": "demo_p1", "text": "Vessel MV-Alpha Retrofit", "start_date": "31-05-2019", "duration": 1460, "progress": 0.5, "open": True, "type": "project" },
                    
                    # Milestones
                    { "id": "m1", "text": "Notice to Proceed", "start_date": "31-05-2019", "duration": 0, "progress": 1, "open": True, "type": "milestone", "parent": "demo_p1" },
                    { "id": "m2", "text": "SRF Notice to Proceed", "start_date": "21-06-2019", "duration": 0, "progress": 1, "open": True, "type": "milestone", "parent": "demo_p1" },
                    { "id": "m3", "text": "Phase 0 Completion", "start_date": "10-07-2022", "duration": 0, "progress": 0, "open": True, "type": "milestone", "parent": "demo_p1" },
                    { "id": "m4", "text": "Phase 1 Completion", "start_date": "26-10-2022", "duration": 0, "progress": 0, "open": True, "type": "milestone", "parent": "demo_p1" },
                    { "id": "m5", "text": "Phase 2 Completion", "start_date": "15-01-2023", "duration": 0, "progress": 0, "open": True, "type": "milestone", "parent": "demo_p1" },
                    { "id": "m6", "text": "Required Phase 0 Pre-Com Completion", "start_date": "20-12-2022", "duration": 0, "progress": 0, "open": True, "type": "milestone", "parent": "demo_p1" },
                    { "id": "m7", "text": "Required Phase 1 Pre-Com Completion", "start_date": "21-03-2023", "duration": 0, "progress": 0, "open": True, "type": "milestone", "parent": "demo_p1" },
                    { "id": "m8", "text": "Required Phase 2 Pre-Com Completion", "start_date": "21-06-2023", "duration": 0, "progress": 0, "open": True, "type": "milestone", "parent": "demo_p1" },

                    # Engineering & Systems
                    { "id": "eng", "text": "Engineering & Systems Execution", "start_date": "01-07-2019", "duration": 1000, "progress": 0.8, "open": True, "type": "project", "parent": "demo_p1" },
                    { "id": "e1", "text": "XT Systems #01 - #02", "start_date": "01-07-2019", "duration": 365, "progress": 1, "parent": "eng" },
                    { "id": "e2", "text": "XT Systems #03 - #06", "start_date": "01-10-2019", "duration": 400, "progress": 1, "parent": "eng" },
                    { "id": "e3", "text": "XT Systems #07 - #10", "start_date": "01-01-2020", "duration": 450, "progress": 0.9, "parent": "eng" },
                    { "id": "e4", "text": "XT Systems #11 - #14", "start_date": "01-04-2020", "duration": 450, "progress": 0.8, "parent": "eng" },
                    { "id": "e5", "text": "XT Systems #15 - #18", "start_date": "01-07-2020", "duration": 450, "progress": 0.6, "parent": "eng" },
                    { "id": "e6", "text": "XT Systems #19 - #20", "start_date": "01-10-2020", "duration": 450, "progress": 0.5, "parent": "eng" },
                    { "id": "e7", "text": "Onshore Controls", "start_date": "01-07-2019", "duration": 300, "progress": 1, "parent": "eng" },
                    { "id": "e8", "text": "CWOR #01 & #02 / IWOCS", "start_date": "01-08-2019", "duration": 450, "progress": 0.9, "parent": "eng" },
                    { "id": "e9", "text": "All Bore Connectors", "start_date": "01-09-2019", "duration": 300, "progress": 0.9, "parent": "eng" },
                    { "id": "e10", "text": "RTM Mounted Controls", "start_date": "01-10-2019", "duration": 400, "progress": 0.8, "parent": "eng" },
                    { "id": "e11", "text": "22\" Connectors", "start_date": "01-11-2019", "duration": 300, "progress": 0.9, "parent": "eng" },

                    # Procurement & Pipeline
                    { "id": "proc", "text": "Procurement & Pipeline Activities", "start_date": "01-09-2019", "duration": 1000, "progress": 0.6, "open": False, "type": "project", "parent": "demo_p1" },
                    { "id": "p1", "text": "Shallow Water Pipe 22\", 8\", 6\"", "start_date": "01-09-2019", "duration": 300, "progress": 1, "parent": "proc" },
                    { "id": "p2", "text": "Deep Water 22\" Pipe (Solitaire)", "start_date": "01-11-2019", "duration": 400, "progress": 0.8, "parent": "proc" },
                    { "id": "p3", "text": "Deep Water 8\", 6\", 4\" & 10\"", "start_date": "01-01-2020", "duration": 400, "progress": 0.7, "parent": "proc" },
                    { "id": "p4", "text": "Deep Water 10\" for Stalking", "start_date": "01-03-2020", "duration": 350, "progress": 0.6, "parent": "proc" },
                    { "id": "p5", "text": "Seamless & Clad Pipe NDE", "start_date": "01-10-2019", "duration": 200, "progress": 1, "parent": "proc" },
                    { "id": "p6", "text": "FJC Qualifications", "start_date": "01-12-2019", "duration": 150, "progress": 0.9, "parent": "proc" },
                    { "id": "p7", "text": "Pipe Transportation", "start_date": "01-05-2020", "duration": 600, "progress": 0.5, "parent": "proc" },

                    # Fabrication
                    { "id": "fab", "text": "Fabrication Activities", "start_date": "01-01-2020", "duration": 800, "progress": 0.7, "open": False, "type": "project", "parent": "demo_p1" },
                    { "id": "f1", "text": "Forgings", "start_date": "01-01-2020", "duration": 400, "progress": 1, "parent": "fab" },
                    { "id": "f2", "text": "PLETs and ILT Fabrication", "start_date": "01-06-2020", "duration": 500, "progress": 0.8, "parent": "fab" },
                    { "id": "f3", "text": "10\" PLETs Fabrication", "start_date": "01-07-2020", "duration": 400, "progress": 0.7, "parent": "fab" },
                    { "id": "f4", "text": "Jumper Pre-Fab / Qualification", "start_date": "01-08-2020", "duration": 600, "progress": 0.4, "parent": "fab" },
                    { "id": "f5", "text": "Gooseneck Fabrication", "start_date": "01-05-2020", "duration": 300, "progress": 0.9, "parent": "fab" },

                    # Offshore Installation
                    { "id": "off", "text": "Offshore / Marine Operations", "start_date": "01-01-2021", "duration": 700, "progress": 0.3, "open": True, "type": "project", "parent": "demo_p1" },
                    { "id": "o1", "text": "Vessel Mobilization (Solitaire / G1200)", "start_date": "01-01-2021", "duration": 60, "progress": 1, "parent": "off" },
                    { "id": "o2", "text": "BHD Dredging", "start_date": "01-02-2021", "duration": 150, "progress": 1, "parent": "off" },
                    { "id": "o3", "text": "Artemis", "start_date": "01-03-2021", "duration": 120, "progress": 0.9, "parent": "off" },
                    { "id": "o4", "text": "Pipe Lay", "start_date": "01-04-2021", "duration": 400, "progress": 0.4, "parent": "off" },
                    { "id": "o5", "text": "Accommodation Setup (S.S M Float)", "start_date": "01-05-2021", "duration": 60, "progress": 0.8, "parent": "off" },
                    { "id": "o6", "text": "Umbilical Installation", "start_date": "01-06-2021", "duration": 300, "progress": 0.2, "parent": "off" },
                    { "id": "o7", "text": "Backfill operations", "start_date": "01-09-2021", "duration": 200, "progress": 0, "parent": "off" },

                    # Commissioning
                    { "id": "com", "text": "Pre-Commissioning & Commissioning", "start_date": "01-06-2022", "duration": 365, "progress": 0.1, "open": True, "type": "project", "parent": "demo_p1" },
                    { "id": "c1", "text": "Umb Pre-Com", "start_date": "01-06-2022", "duration": 90, "progress": 0.5, "parent": "com" },
                    { "id": "c2", "text": "Pipe Pre-Com", "start_date": "01-08-2022", "duration": 120, "progress": 0, "parent": "com" },
                    { "id": "c3", "text": "Fuel Gas Testing", "start_date": "01-10-2022", "duration": 60, "progress": 0, "parent": "com" },
                    { "id": "c4", "text": "System Handover", "start_date": "01-12-2022", "duration": 180, "progress": 0, "parent": "com" }
                ],
                "links": [
                    { "id": 1, "source": "m1", "target": "eng", "type": "0" },
                    { "id": 2, "source": "e1", "target": "fab", "type": "0" },
                    { "id": 3, "source": "f2", "target": "off", "type": "0" },
                    { "id": 4, "source": "o4", "target": "com", "type": "0" },
                    { "id": 5, "source": "c4", "target": "m5", "type": "0" },
                    { "id": 6, "source": "p6", "target": "fab", "type": "0" },
                    { "id": 7, "source": "f4", "target": "p7", "type": "0" },
                    { "id": 8, "source": "p7", "target": "off", "type": "0" },
                    { "id": 9, "source": "o2", "target": "o4", "type": "0" }
                ]
            }
            
            # For daily timeline config, keep demo durations in days directly (do not divide by 7)
            return demo_data

        from sqlalchemy.orm import joinedload
        project = db.query(Project).options(
            joinedload(Project.tasks).joinedload(Task.subtasks)
        ).filter(Project.project_id == project_id).first()
        
        if not project:
            raise HTTPException(status_code=404, detail="Project not found")

        gantt_data = []

        # Determine project start and end dates from tasks to ensure they match
        valid_tasks = [t for t in project.tasks if t.planned_start_date and t.planned_end_date]
        if valid_tasks:
            from decimal import Decimal
            proj_start = min(t.planned_start_date for t in valid_tasks)
            proj_end = max(t.planned_end_date for t in valid_tasks)
            
            # Heal database project record if out of sync
            proj_work_days = WeekendHelper.count_work_days(proj_start, proj_end)
            expected_weeks = WeekendHelper.working_days_to_weeks(proj_work_days)
            if project.planned_start_date != proj_start or project.duration_weeks != expected_weeks:
                project.planned_start_date = proj_start
                project.duration_weeks = expected_weeks
                project.budget_at_completion = sum(Decimal(str(t.planned_cost or 0)) for t in project.tasks)
                db.commit()
                db.refresh(project)
                
            proj_start_str = proj_start.strftime("%d-%m-%Y")
            # Calculate duration in work days (excluding weekends)
            proj_duration_days = max(proj_work_days, 1)
        else:
            proj_start = project.planned_start_date or datetime.now().date()
            proj_start_str = proj_start.strftime("%d-%m-%Y")
            proj_duration_days = WeekendHelper.weeks_to_working_days(project.duration_weeks or 5)
            proj_end = WeekendHelper.add_work_days(proj_start, proj_duration_days)

        proj_actual_cost = float(project.actual_cost) if project.actual_cost is not None else sum(float(t.actual_cost or 0) for t in project.tasks)
        
        valid_actual_tasks = [t for t in project.tasks if t.actual_start_date and t.actual_end_date]
        if valid_actual_tasks:
            proj_act_start = min(t.actual_start_date for t in valid_actual_tasks)
            proj_act_end = max(t.actual_end_date for t in valid_actual_tasks)
            proj_act_start_str = proj_act_start.strftime("%d-%m-%Y")
            proj_act_end_str = proj_act_end.strftime("%d-%m-%Y")
            proj_act_duration = WeekendHelper.count_work_days(proj_act_start, proj_act_end)
        else:
            proj_act_start = project.actual_start_date or (proj_start if 'proj_start' in locals() else datetime.now().date())
            proj_act_start_str = proj_act_start.strftime("%d-%m-%Y") if proj_act_start else ""
            proj_act_duration = WeekendHelper.weeks_to_working_days(project.actual_duration or 5)
            proj_act_end = WeekendHelper.add_work_days(proj_act_start, proj_act_duration) if proj_act_start else None
            proj_act_end_str = proj_act_end.strftime("%d-%m-%Y") if proj_act_end else ""

        proj_planned_start_str = proj_start_str
        proj_planned_end_str = proj_end.strftime("%d-%m-%Y") if proj_end else ""
        proj_planned_duration = proj_duration_days

        proj_crew = project.crew_size or sum(t.planned_crew for t in project.tasks if t.planned_crew) or 1
        proj_act_crew = project.crew_size or sum(t.actual_crew for t in project.tasks if t.actual_crew) or proj_crew

        # Compute project progress on-the-fly as the average of its tasks' progress
        proj_progress = 0.0
        if project.tasks:
            task_progress_list = []
            for t in project.tasks:
                if t.subtasks:
                    sub_progs = [float(s.progress_percent or 0.0) for s in t.subtasks]
                    task_prog = sum(sub_progs) / len(t.subtasks) if t.subtasks else 0.0
                else:
                    task_prog = float(t.progress_percent or 0.0)
                task_progress_list.append(task_prog)
            proj_progress = (sum(task_progress_list) / len(project.tasks)) / 100.0 if project.tasks else 0.0

        gantt_data.append({
            "id": str(project.project_id),
            "text": project.project_title,
            "start_date": proj_act_start_str,
            "duration": proj_act_duration,
            "progress": proj_progress,
            "open": True,
            "type": "project",
            "status": project.status or "In Progress",
            "unscheduled": not bool(proj_act_start_str),
            "Actual_start_date": proj_act_start_str,
            "Actual_end_date": proj_act_end_str,
            "Actual_duration": proj_act_duration,
            "Actual_crew": proj_act_crew,
            "Actual_cost": proj_actual_cost,
            "cost": proj_actual_cost,
            "planned_start_date": proj_planned_start_str,
            "planned_end_date": proj_planned_end_str,
            "Planned_duration": proj_planned_duration,
            "Planned_crew": proj_crew,
            "Planned_cost": float(project.budget_at_completion or 0),
            # Legacy casing keys
            "actual_crew": proj_act_crew,
            "planned_crew": proj_crew,
            "actual_cost": proj_actual_cost,
            "planned_cost": float(project.budget_at_completion or 0),
            "actual_end_date": proj_act_end_str,
            "planned_duration": proj_planned_duration
        })
        
        for task in project.tasks:
            # 1. Start/End date strictly from actual fields
            t_act_start = task.actual_start_date
            t_act_end = task.actual_end_date
            
            if t_act_start:
                task_start_str = t_act_start.strftime("%d-%m-%Y")
            else:
                task_start_str = ""
            
            # 2. Calculate duration in working days (excluding weekends)
            if t_act_start and t_act_end:
                t_dur = max(WeekendHelper.count_work_days(t_act_start, t_act_end), 1)
            else:
                t_dur = "" # empty default

            if task.subtasks:
                sub_progs = [float(s.progress_percent or 0.0) for s in task.subtasks]
                task_progress = (sum(sub_progs) / len(task.subtasks)) / 100.0 if task.subtasks else 0.0
            else:
                task_progress = float(task.progress_percent or 0.0) / 100.0

            gantt_data.append({
                "id": str(task.task_id),
                "text": task.task_name,
                "start_date": task_start_str,
                "duration": t_dur,
                "cost": float(task.actual_cost or 0),
                "parent": str(project.project_id),
                "progress": task_progress,
                "open": True,
                "type": "task",
                "status": task.status or "In Progress",
                "unscheduled": not bool(task_start_str),
                "Actual_start_date": task.actual_start_date.strftime("%d-%m-%Y") if task.actual_start_date else "",
                "Actual_end_date": task.actual_end_date.strftime("%d-%m-%Y") if task.actual_end_date else "",
                "Actual_duration": t_dur if isinstance(t_dur, int) else 0,
                "Actual_crew": task.actual_crew,
                "Actual_cost": float(task.actual_cost or 0),
                "planned_start_date": task.planned_start_date.strftime("%d-%m-%Y") if task.planned_start_date else "",
                "planned_end_date": task.planned_end_date.strftime("%d-%m-%Y") if task.planned_end_date else "",
                "Planned_duration": WeekendHelper.count_work_days(task.planned_start_date, task.planned_end_date) if task.planned_start_date and task.planned_end_date else 0,
                "Planned_crew": task.planned_crew,
                "Planned_cost": float(task.planned_cost or 0),
                # Legacy keys for backward compatibility
                "actual_end_date": task.actual_end_date.strftime("%d-%m-%Y") if task.actual_end_date else "",
                "planned_end_date": task.planned_end_date.strftime("%d-%m-%Y") if task.planned_end_date else "",
                "planned_start_date": task.planned_start_date.strftime("%d-%m-%Y") if task.planned_start_date else "",
                "planned_duration": WeekendHelper.count_work_days(task.planned_start_date, task.planned_end_date) if task.planned_start_date and task.planned_end_date else 0,
                "planned_crew": task.planned_crew,
                "actual_crew": task.actual_crew,
                "planned_cost": float(task.planned_cost or 0),
                "actual_cost": float(task.actual_cost or 0)
            })
            
            for sub in task.subtasks:
                sub_start_obj = sub.actual_start_date
                sub_end_obj = sub.actual_end_date
                
                if sub_start_obj:
                    sub_start_str = sub_start_obj.strftime("%d-%m-%Y")
                else:
                    sub_start_str = ""
                
                # Calculate subtask duration in working days (excluding weekends)
                if sub_start_obj and sub_end_obj:
                    s_dur = max(WeekendHelper.count_work_days(sub_start_obj, sub_end_obj), 1)
                else:
                    s_dur = "" # empty default

                gantt_data.append({
                    "id": str(sub.subtask_id),
                    "text": sub.subtask_name,
                    "start_date": sub_start_str,
                    "duration": s_dur,
                    "cost": float(sub.actual_cost or 0),
                    "parent": str(task.task_id),
                    "progress": float(sub.progress_percent or 0) / 100.0,
                    "open": True,
                    "type": "task",
                    "status": sub.status or "In Progress",
                    "unscheduled": not bool(sub_start_str),
                    "Actual_start_date": sub.actual_start_date.strftime("%d-%m-%Y") if sub.actual_start_date else "",
                    "Actual_end_date": sub.actual_end_date.strftime("%d-%m-%Y") if sub.actual_end_date else "",
                    "Actual_duration": s_dur if isinstance(s_dur, int) else 0,
                    "Actual_crew": sub.actual_crew,
                    "Actual_cost": float(sub.actual_cost or 0),
                    "planned_start_date": sub.planned_start_date.strftime("%d-%m-%Y") if sub.planned_start_date else "",
                    "planned_end_date": sub.planned_end_date.strftime("%d-%m-%Y") if sub.planned_end_date else "",
                    "Planned_duration": WeekendHelper.count_work_days(sub.planned_start_date, sub.planned_end_date) if sub.planned_start_date and sub.planned_end_date else 0,
                    "Planned_crew": sub.planned_crew,
                    "Planned_cost": float(sub.planned_cost or 0),
                    # Legacy keys for backward compatibility
                    "actual_end_date": sub.actual_end_date.strftime("%d-%m-%Y") if sub.actual_end_date else "",
                    "planned_end_date": sub.planned_end_date.strftime("%d-%m-%Y") if sub.planned_end_date else "",
                    "planned_start_date": sub.planned_start_date.strftime("%d-%m-%Y") if sub.planned_start_date else "",
                    "planned_duration": WeekendHelper.count_work_days(sub.planned_start_date, sub.planned_end_date) if sub.planned_start_date and sub.planned_end_date else 0,
                    "planned_crew": sub.planned_crew,
                    "actual_crew": sub.actual_crew,
                    "planned_cost": float(sub.planned_cost or 0),
                    "actual_cost": float(sub.actual_cost or 0)
                })
                
        task_links = db.query(TaskLink).filter(TaskLink.project_id == project_id).all()
        gantt_links = [{
            "id": str(link.link_id),
            "source": link.source_id,
            "target": link.target_id,
            "type": link.link_type
        } for link in task_links]
        
        return {"data": gantt_data, "links": gantt_links}
        
    except HTTPException:
        raise
    except Exception as e:
        logger.error(f"Error fetching gantt data: {str(e)}")
        raise HTTPException(status_code=500, detail="Internal Server Error")

def propagate_project_dates(project, db: Session):
    from decimal import Decimal
    # 1. Planned metrics propagation
    valid_tasks = [t for t in project.tasks if t.planned_start_date and t.planned_end_date]
    if valid_tasks:
        min_start = min(t.planned_start_date for t in valid_tasks)
        max_end = max(t.planned_end_date for t in valid_tasks)
        
        project.planned_start_date = min_start
        
        proj_work_days = WeekendHelper.count_work_days(min_start, max_end)
        project.planned_duration = WeekendHelper.working_days_to_weeks(proj_work_days)
        
        # Recalculate project budget (planned cost)
        project.planned_cost = sum(Decimal(str(t.planned_cost or 0)) for t in project.tasks)
    else:
        project.planned_start_date = None
        project.planned_duration = 0
        project.planned_cost = Decimal('0.00')

    # 2. Actual metrics propagation
    valid_actual_tasks = [t for t in project.tasks if t.actual_start_date and t.actual_end_date]
    if valid_actual_tasks:
        min_act_start = min(t.actual_start_date for t in valid_actual_tasks)
        max_act_end = max(t.actual_end_date for t in valid_actual_tasks)
        
        project.actual_start_date = min_act_start
        
        proj_act_work_days = WeekendHelper.count_work_days(min_act_start, max_act_end)
        project.actual_duration = WeekendHelper.working_days_to_weeks(proj_act_work_days)
        
        # Recalculate project actual cost
        project.actual_cost = sum(Decimal(str(t.actual_cost or 0)) for t in project.tasks)
    else:
        project.actual_start_date = None
        project.actual_duration = 0
        project.actual_cost = Decimal('0.00')

def propagate_dates_upward(entity, db: Session):
    from decimal import Decimal
    if isinstance(entity, Subtask):
        parent_task = entity.task
        if parent_task:
            # Recalculate progress as average of subtasks
            if parent_task.subtasks:
                sub_progs = [float(s.progress_percent or 0.0) for s in parent_task.subtasks]
                parent_task.progress_percent = round(sum(sub_progs) / len(parent_task.subtasks), 2)

            # 1. Planned dates propagation
            valid_planned_subtasks = [s for s in parent_task.subtasks if s.planned_start_date and s.planned_end_date]
            if valid_planned_subtasks:
                min_p_start = min(s.planned_start_date for s in valid_planned_subtasks)
                max_p_end = max(s.planned_end_date for s in valid_planned_subtasks)
                parent_task.planned_start_date = min_p_start
                parent_task.planned_end_date = max_p_end
                
                # Rule 2: Task Planned Budget = Sum of its child subtasks' planned costs.
                parent_task.planned_cost = sum(s.planned_cost or Decimal('0.00') for s in parent_task.subtasks)
            else:
                parent_task.planned_start_date = None
                parent_task.planned_end_date = None
                parent_task.planned_cost = Decimal('0.00')
            
            # 2. Actual dates propagation
            valid_actual_subtasks = [s for s in parent_task.subtasks if s.actual_start_date and s.actual_end_date]
            if valid_actual_subtasks:
                min_a_start = min(s.actual_start_date for s in valid_actual_subtasks)
                max_a_end = max(s.actual_end_date for s in valid_actual_subtasks)
                parent_task.actual_start_date = min_a_start
                parent_task.actual_end_date = max_a_end
                
                # Parent task actual cost is sum of subtask actual costs.
                parent_task.actual_cost = sum(s.actual_cost or Decimal('0.00') for s in parent_task.subtasks)
            else:
                parent_task.actual_start_date = None
                parent_task.actual_end_date = None
                parent_task.actual_cost = Decimal('0.00')
                
            # Propagate from parent task to project
            if parent_task.project:
                propagate_project_dates(parent_task.project, db)
                    
    elif isinstance(entity, Task):
        # Recalculate progress as average of subtasks
        if entity.subtasks:
            sub_progs = [float(s.progress_percent or 0.0) for s in entity.subtasks]
            entity.progress_percent = round(sum(sub_progs) / len(entity.subtasks), 2)

        valid_planned_subtasks = [s for s in entity.subtasks if s.planned_start_date and s.planned_end_date]
        if valid_planned_subtasks:
            min_p_start = min(s.planned_start_date for s in valid_planned_subtasks)
            max_p_end = max(s.planned_end_date for s in valid_planned_subtasks)
            entity.planned_start_date = min_p_start
            entity.planned_end_date = max_p_end
            
            # Rule 2: Task Planned Budget = Sum of its child subtasks' planned costs.
            entity.planned_cost = sum(s.planned_cost or Decimal('0.00') for s in entity.subtasks)
        else:
            if entity.subtasks:
                entity.planned_start_date = None
                entity.planned_end_date = None
                entity.planned_cost = Decimal('0.00')
            
        valid_actual_subtasks = [s for s in entity.subtasks if s.actual_start_date and s.actual_end_date]
        if valid_actual_subtasks:
            min_a_start = min(s.actual_start_date for s in valid_actual_subtasks)
            max_a_end = max(s.actual_end_date for s in valid_actual_subtasks)
            entity.actual_start_date = min_a_start
            entity.actual_end_date = max_a_end
            
            # Parent task actual cost is sum of subtask actual costs.
            entity.actual_cost = sum(s.actual_cost or Decimal('0.00') for s in entity.subtasks)
        else:
            if entity.subtasks:
                entity.actual_start_date = None
                entity.actual_end_date = None
                entity.actual_cost = Decimal('0.00')
            
        if entity.project:
            propagate_project_dates(entity.project, db)

@router.post("/{project_id}/gantt/tasks")
def create_gantt_task(project_id: str, payload: GanttTaskUpdate, db: Session = Depends(get_db)):
    try:
        from datetime import datetime
        from core.date_helper import WeekendHelper
        from core.budget_helper import BudgetHelper
        from decimal import Decimal
        
        # 1. Parse start date
        start_obj = None
        payload_start_date = payload.Actual_start_date or payload.start_date or payload.planned_start_date
        if payload_start_date:
            try:
                start_date_str = payload_start_date.split(" ")[0]
                for fmt in ("%d-%m-%Y", "%Y-%m-%d", "%Y-%m-%dT%H:%M:%S"):
                    try:
                        start_obj = datetime.strptime(start_date_str, fmt).date()
                        break
                    except ValueError:
                        continue
            except Exception as e:
                logger.error(f"Error parsing task start_date: {str(e)}")
                
        if not start_obj:
            project = db.query(Project).filter(Project.project_id == project_id).first()
            if project and project.planned_start_date:
                start_obj = project.planned_start_date
            else:
                start_obj = datetime.now().date()
            start_obj = WeekendHelper.roll_to_next_work_day(start_obj)
            
        # 2. Calculate end date based on duration
        duration = payload.Actual_duration if payload.Actual_duration is not None else (payload.duration if payload.duration is not None else (payload.Planned_duration or 1))
        end_obj = WeekendHelper.add_work_days(start_obj, duration)
        
        # 3. Determine if it is a Task or a Subtask
        is_main_task = str(payload.parent) == str(project_id)
        
        new_id = uuid.uuid4()
        progress_pct = round((payload.progress or 0.0) * 100, 2)
        
        planned_crew = payload.Planned_crew if payload.Planned_crew is not None else (payload.planned_crew if payload.planned_crew is not None else 1)
        actual_crew = payload.Actual_crew if payload.Actual_crew is not None else (payload.actual_crew if payload.actual_crew is not None else planned_crew)
        
        calculated_p_cost = Decimal(str(BudgetHelper.calculate_budget(start_obj, end_obj, planned_crew)))
        calculated_a_cost = Decimal(str(BudgetHelper.calculate_budget(start_obj, end_obj, actual_crew)))
        
        if is_main_task:
            task_code = f"TSK-{str(uuid.uuid4())[:8].upper()}"
            new_task = Task(
                task_id=new_id,
                task_code=task_code,
                project_id=project_id,
                task_name=payload.text or "New Task",
                planned_start_date=start_obj,
                planned_end_date=end_obj,
                actual_start_date=start_obj,
                actual_end_date=end_obj,
                planned_cost=calculated_p_cost,
                actual_cost=calculated_a_cost,
                progress_percent=progress_pct,
                status="To Do",
                priority="Medium",
                planned_crew=planned_crew,
                actual_crew=actual_crew,
                order_index=1
            )
            db.add(new_task)
            entity = new_task
        else:
            parent_task = db.query(Task).filter(Task.task_id == payload.parent).first()
            if not parent_task:
                raise HTTPException(status_code=404, detail="Parent task not found")
                
            sub_code = f"SUB-{str(uuid.uuid4())[:8].upper()}"
            new_subtask = Subtask(
                subtask_id=new_id,
                subtask_code=sub_code,
                task_id=parent_task.task_id,
                subtask_name=payload.text or "New Subtask",
                planned_start_date=start_obj,
                planned_end_date=end_obj,
                actual_start_date=start_obj,
                actual_end_date=end_obj,
                planned_cost=calculated_p_cost,
                actual_cost=calculated_a_cost,
                progress_percent=progress_pct,
                status="To Do",
                planned_crew=planned_crew,
                actual_crew=actual_crew,
                order_index=1
            )
            db.add(new_subtask)
            entity = new_subtask
            
        db.flush()
        
        # Propagate dates upward
        propagate_dates_upward(entity, db)
        
        db.commit()
        return {"action": "inserted", "tid": str(new_id)}
    except HTTPException:
        db.rollback()
        raise
    except Exception as e:
        db.rollback()
        logger.error(f"Error creating task: {str(e)}")
        raise HTTPException(status_code=500, detail="Internal Server Error")

@router.put("/{project_id}/gantt/tasks/{task_id}")
def update_gantt_task(project_id: str, task_id: str, payload: GanttTaskUpdate, db: Session = Depends(get_db)):
    try:
        from datetime import datetime, timedelta
        from core.budget_helper import BudgetHelper
        
        # Determine if it's a project, task, or subtask.
        if task_id == project_id:
            project = db.query(Project).filter(Project.project_id == project_id).first()
            if not project:
                raise HTTPException(status_code=404, detail="Project not found")
            entity = project
        else:
            task = db.query(Task).filter(Task.task_id == task_id).first()
            if not task:
                subtask = db.query(Subtask).filter(Subtask.subtask_id == task_id).first()
                if not subtask:
                    raise HTTPException(status_code=404, detail="Task or Subtask not found")
                entity = subtask
            else:
                entity = task
            
        # Determine if the entity is a Task that has subtasks in the database
        is_task_with_subtasks = False
        if isinstance(entity, Task):
            is_task_with_subtasks = len(entity.subtasks) > 0

        dates_changed = False
        needs_propagation = False
        
        # Helper to parse dates
        def parse_date_str(date_str: str):
            if not date_str:
                return None
            try:
                date_str_clean = date_str.split(" ")[0]
                for fmt in ("%d-%m-%Y", "%Y-%m-%d", "%Y-%m-%dT%H:%M:%S"):
                    try:
                        return datetime.strptime(date_str_clean, fmt).date()
                    except ValueError:
                        continue
            except Exception:
                pass
            return None

        # Only parse and apply payload date changes if it is NOT a task with subtasks
        if not is_task_with_subtasks:
            if isinstance(entity, Project):
                # For Project, map start_date to planned_start_date
                proj_start_str = payload.Actual_start_date or payload.start_date or payload.planned_start_date
                proj_start_obj = parse_date_str(proj_start_str)
                if proj_start_obj and entity.planned_start_date != proj_start_obj:
                    entity.planned_start_date = proj_start_obj
                    dates_changed = True
                
                proj_duration = payload.Actual_duration if payload.Actual_duration is not None else payload.duration
                if proj_duration is not None:
                    new_weeks = WeekendHelper.working_days_to_weeks(proj_duration)
                    if entity.duration_weeks != new_weeks:
                        entity.duration_weeks = new_weeks
                        dates_changed = True
            else:
                # 1. ACTUAL Dates for Task / Subtask
                act_start_str = payload.Actual_start_date or payload.start_date
                act_duration = payload.Actual_duration if payload.Actual_duration is not None else payload.duration
                
                act_start_obj = parse_date_str(act_start_str)
                if act_start_obj:
                    if entity.actual_start_date != act_start_obj:
                        entity.actual_start_date = act_start_obj
                        dates_changed = True
                    
                    if act_duration is not None:
                        act_end_obj = WeekendHelper.add_work_days(act_start_obj, act_duration)
                        if entity.actual_end_date != act_end_obj:
                            entity.actual_end_date = act_end_obj
                            dates_changed = True
                
                # 2. PLANNED Dates for Task / Subtask
                pl_start_str = payload.planned_start_date
                pl_duration = payload.Planned_duration
                
                pl_start_obj = parse_date_str(pl_start_str)
                if pl_start_obj:
                    if entity.planned_start_date != pl_start_obj:
                        entity.planned_start_date = pl_start_obj
                        dates_changed = True
                    
                    if pl_duration is not None:
                        pl_end_obj = WeekendHelper.add_work_days(pl_start_obj, pl_duration)
                        if entity.planned_end_date != pl_end_obj:
                            entity.planned_end_date = pl_end_obj
                            dates_changed = True

        if payload.status is not None:
            entity.status = payload.status
            if payload.status == "Not Started" and isinstance(entity, (Task, Subtask)):
                entity.actual_start_date = None
                entity.actual_end_date = None
                entity.planned_start_date = None
                entity.planned_end_date = None
                entity.planned_cost = 0.0
                entity.actual_cost = 0.0
                entity.progress_percent = 0.0
                entity.planned_crew = 0
                entity.actual_crew = 0
                dates_changed = True

        if payload.progress is not None and isinstance(entity, (Task, Subtask)):
            entity.progress_percent = round(payload.progress * 100, 2)
            needs_propagation = True
            
        if payload.text is not None:
            if isinstance(entity, Task):
                entity.task_name = payload.text
            elif isinstance(entity, Subtask):
                entity.subtask_name = payload.text
            elif isinstance(entity, Project):
                entity.project_title = payload.text
                
        needs_propagation = needs_propagation or is_task_with_subtasks
            
        if isinstance(entity, (Task, Subtask)):
            crew_changed = False
            planned_crew_val = payload.Planned_crew if payload.Planned_crew is not None else payload.planned_crew
            if planned_crew_val is not None and entity.planned_crew != planned_crew_val:
                entity.planned_crew = planned_crew_val
                crew_changed = True
            actual_crew_val = payload.Actual_crew if payload.Actual_crew is not None else payload.actual_crew
            if actual_crew_val is not None and entity.actual_crew != actual_crew_val:
                entity.actual_crew = actual_crew_val
                crew_changed = True
                
            from decimal import Decimal
            
            # Recalculate costs if dates or crew changed
            if dates_changed or crew_changed:
                calc_p_cost = BudgetHelper.calculate_budget(
                    entity.planned_start_date,
                    entity.planned_end_date,
                    entity.planned_crew
                )
                entity.planned_cost = Decimal(str(calc_p_cost))
                
                calc_a_cost = BudgetHelper.calculate_budget(
                    entity.actual_start_date,
                    entity.actual_end_date,
                    entity.actual_crew
                )
                entity.actual_cost = Decimal(str(calc_a_cost))
                needs_propagation = True
            else:
                # Explicit cost values or fallback
                planned_cost_val = payload.Planned_cost if payload.Planned_cost is not None else payload.planned_cost
                if planned_cost_val is not None:
                    entity.planned_cost = Decimal(str(planned_cost_val))
                    needs_propagation = True
                actual_cost_val = payload.Actual_cost if payload.Actual_cost is not None else payload.actual_cost
                if actual_cost_val is not None:
                    entity.actual_cost = Decimal(str(actual_cost_val))
                    needs_propagation = True
                
            if dates_changed:
                needs_propagation = True
                
            if needs_propagation:
                propagate_dates_upward(entity, db)
                
        db.commit()
        return {"action": "updated"}
    except Exception as e:
        db.rollback()
        logger.error(f"Error updating task: {str(e)}")
        raise HTTPException(status_code=500, detail="Internal Server Error")

@router.get("/{project_id}/evm")
def get_project_evm_data(project_id: str, db: Session = Depends(get_db)):
    """
    Returns EVM timeseries data (Weekly, Monthly, Yearly) for a specific project,
    accompanied by deterministic metrics and structured AI insights.
    """
    # 1. Handle Demo Project
    if project_id == "demo-maritime-project-001":
        return {
            "Weekly": {
                "labels": ["W1", "W2", "W3", "W4", "W5", "W6", "W7", "W8", "W9"],
                "budgeted_cost": [10000.0, 20000.0, 30000.0, 55000.0, 85000.0, 120000.0, 130000.0, 140000.0, 150000.0],
                "actual_cost": [5000.0, 10000.0, 25000.0, 75000.0, 120000.0, 155000.0, 170000.0, None, None],
                "predicted_cost": [5000.0, 10000.0, 25000.0, 75000.0, 120000.0, 155000.0, 170000.0, 205000.0, 240000.0],
                "trend_color": "red",
                "additional_workers_needed": 4,
                "suggestion": "The project is currently over budget by €20k and behind schedule (SPI = 0.69). To complete the retrofit of MV-Alpha on schedule (by W9), it is highly recommended to onboard 4 additional crew members immediately to accelerate the piping and mechanical installation tasks."
            },
            "Monthly": {
                "labels": ["Month 1", "Month 2", "Month 3"],
                "budgeted_cost": [85000.0, 150000.0, 150000.0],
                "actual_cost": [120000.0, 170000.0, None],
                "predicted_cost": [120000.0, 170000.0, 240000.0],
                "trend_color": "red",
                "additional_workers_needed": 4,
                "suggestion": "The project is currently over budget by €20k and behind schedule (SPI = 0.69). To complete the retrofit of MV-Alpha on schedule, it is highly recommended to onboard 4 additional crew members immediately to accelerate the piping and mechanical installation tasks."
            },
            "Yearly": {
                "labels": ["2025", "2026"],
                "budgeted_cost": [150000.0, 150000.0],
                "actual_cost": [170000.0, None],
                "predicted_cost": [170000.0, 240000.0],
                "trend_color": "red",
                "additional_workers_needed": 4,
                "suggestion": "The project is currently over budget by €20k and behind schedule (SPI = 0.69). To complete the retrofit of MV-Alpha on schedule, it is highly recommended to onboard 4 additional crew members immediately to accelerate the piping and mechanical installation tasks."
            },
            "metrics": {
                "bac": 150000.0, "pv": 130000.0, "ev": 90000.0, "ac": 170000.0,
                "cpi": 0.53, "spi": 0.69, "cv": -80000.0, "sv": -40000.0,
                "eac": 283000.0, "etc": 113000.0, "vac": -133000.0,
                "predictedCompletionDate": "2026-08-09",
                "predictedCompletionCost": 283000.0,
                "remainingBudget": 0.0,
                "remainingDurationDays": 14,
                "remainingWorkHours": 320.0,
                "currentCrewUtilization": 85.0,
                "forecastDelayDays": 6,
                "criticalPathCount": 2,
                "tasksNotStartedCount": 1,
                "blockedTasksCount": 1,
                "delayedTasksCount": 2
            },
            "ai_analysis": {
                "forecast": {
                    "predictedCompletionCost": 283000.0,
                    "predictedCompletionDate": "2026-08-09",
                    "confidence": 0.89
                },
                "budget": {
                    "status": "Over Budget",
                    "expectedVariance": -133000.0,
                    "risk": "High"
                },
                "schedule": {
                    "status": "Delayed",
                    "expectedDelayDays": 6,
                    "criticalTasks": ["Hull Inspection", "Pipe Installation"]
                },
                "resources": {
                    "currentCrew": 10,
                    "recommendedCrew": 14,
                    "additionalCrewNeeded": 4
                },
                "recommendations": {
                    "resourceOptimization": [
                        "Add 4 fitters/welders to Pipe Installation immediately.",
                        "Reassign 2 idle electricians from Painting to Hull Inspection."
                    ],
                    "scheduleRecovery": [
                        "Run structural weld tests in parallel with mechanical pipe checks.",
                        "Fast-track paint procurement to avoid surface preparation delay."
                    ],
                    "budgetOptimization": [
                        "Minimize overtime on structural painting packages.",
                        "Delay non-critical cosmetic outfitting tasks to protect remaining contingency."
                    ],
                    "riskMitigation": [
                        "High risk: Weld inspection has not started and holds up the critical path.",
                        "Weather alert: Late season rain may impact outdoor drydock painting."
                    ]
                }
            }
        }

    # 2. Call the EVM Service
    from services.evm_service import EVMService
    return EVMService.calculate_project_evm(project_id, db)


@router.post("/{project_id}/gantt/links")
def create_gantt_link(project_id: str, payload: GanttLinkCreate, db: Session = Depends(get_db)):
    try:
        link = TaskLink(
            link_id=str(payload.id),
            project_id=project_id,
            source_id=str(payload.source),
            target_id=str(payload.target),
            link_type=str(payload.type)
        )
        db.add(link)
        db.commit()
        return {"action": "inserted", "tid": payload.id}
    except Exception as e:
        db.rollback()
        logger.error(f"Error creating link: {str(e)}")
        raise HTTPException(status_code=500, detail="Internal Server Error")

@router.put("/{project_id}/gantt/links/{link_id}")
def update_gantt_link(project_id: str, link_id: str, payload: GanttLinkUpdate, db: Session = Depends(get_db)):
    try:
        link = db.query(TaskLink).filter(TaskLink.link_id == link_id, TaskLink.project_id == project_id).first()
        if link:
            link.source_id = str(payload.source)
            link.target_id = str(payload.target)
            link.link_type = str(payload.type)
            db.commit()
        return {"action": "updated"}
    except Exception as e:
        db.rollback()
        logger.error(f"Error updating link: {str(e)}")
        raise HTTPException(status_code=500, detail="Internal Server Error")

@router.delete("/{project_id}/gantt/links/{link_id}")
def delete_gantt_link(project_id: str, link_id: str, db: Session = Depends(get_db)):
    try:
        db.query(TaskLink).filter(TaskLink.link_id == link_id, TaskLink.project_id == project_id).delete()
        db.commit()
        return {"action": "deleted"}
    except Exception as e:
        db.rollback()
        logger.error(f"Error deleting link: {str(e)}")
        raise HTTPException(status_code=500, detail="Internal Server Error")

@router.delete("/{project_id}/gantt/tasks/{task_id}")
def delete_gantt_task(project_id: str, task_id: str, db: Session = Depends(get_db)):
    try:
        from decimal import Decimal
        # Determine if it's a task or subtask.
        task = db.query(Task).filter(Task.task_id == task_id).first()
        if not task:
            subtask = db.query(Subtask).filter(Subtask.subtask_id == task_id).first()
            if not subtask:
                raise HTTPException(status_code=404, detail="Task or Subtask not found")
            
            parent_task = subtask.task
            db.delete(subtask)
            db.flush()
            
            if parent_task:
                # Update parent task dates and costs
                valid_subtasks = [s for s in parent_task.subtasks if s.planned_start_date and s.planned_end_date]
                if valid_subtasks:
                    min_start = min(s.planned_start_date for s in valid_subtasks)
                    max_end = max(s.planned_end_date for s in valid_subtasks)
                    parent_task.planned_start_date = min_start
                    parent_task.actual_start_date = min_start
                    parent_task.planned_end_date = max_end
                    parent_task.actual_end_date = max_end
                    
                    # Recalculate cost of parent task as sum of subtasks' costs
                    parent_task.planned_cost = sum(s.planned_cost or Decimal('0.00') for s in parent_task.subtasks)
                    parent_task.actual_cost = sum(s.actual_cost or Decimal('0.00') for s in parent_task.subtasks)
                else:
                    # No subtasks left, reset parent task dates
                    parent_task.planned_start_date = None
                    parent_task.actual_start_date = None
                    parent_task.planned_end_date = None
                    parent_task.actual_end_date = None
                    parent_task.planned_cost = Decimal('0.00')
                    parent_task.actual_cost = Decimal('0.00')
                
                if parent_task.project:
                    propagate_project_dates(parent_task.project, db)
        else:
            project = task.project
            db.delete(task)
            db.flush()
            
            if project:
                propagate_project_dates(project, db)
                
        db.commit()
        return {"action": "deleted"}
    except Exception as e:
        db.rollback()
        logger.error(f"Error deleting task: {str(e)}")
        raise HTTPException(status_code=500, detail="Internal Server Error")
